import os
import re
from datetime import datetime
from typing import Any, Dict, List, Optional

from docx import Document
from nltk.tokenize import sent_tokenize
from pdf2image import convert_from_path
from pptx import Presentation
import pytesseract

from policymind.core.config import Settings
from policymind.models.schemas import DocumentMetadata, DocumentType

try:
    import fitz

    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import nltk

    nltk.data.find("tokenizers/punkt")
    nltk.data.find("tokenizers/punkt_tab")
except (ImportError, LookupError, OSError):
    import nltk

    nltk.download("punkt", quiet=True)
    nltk.download("punkt_tab", quiet=True)


class DocumentProcessor:
    def __init__(self, settings: Settings):
        self.settings = settings
        self.chunk_size = settings.CHUNK_SIZE
        self.chunk_overlap = settings.CHUNK_OVERLAP

    def process_document(self, file_path: str, document_id: str) -> Dict[str, Any]:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found at path: {file_path}")

        text_content, metadata = self._extract_text_and_metadata(file_path, document_id)
        chunks = self._create_chunks(text_content, metadata)
        return {
            "document_id": document_id,
            "chunks": chunks,
            "metadata": metadata,
            "total_chunks": len(chunks),
        }

    def _extract_text_and_metadata(self, file_path: str, document_id: str) -> tuple[str, DocumentMetadata]:
        ext = os.path.splitext(file_path)[1].lower()
        processor_map = {
            ".pdf": self._process_pdf,
            ".docx": self._process_docx,
            ".txt": self._process_txt,
            ".pptx": self._process_pptx,
        }
        if ext not in processor_map:
            raise ValueError(f"Unsupported file format: {ext}")

        text, page_count = processor_map[ext](file_path)
        cleaned_text = self._clean_text(text)
        metadata = self._create_metadata(document_id, cleaned_text, file_path, page_count)
        return cleaned_text, metadata

    def _process_pdf(self, file_path: str) -> tuple[str, int]:
        text = ""
        page_count = 0

        if PYMUPDF_AVAILABLE:
            try:
                with fitz.open(file_path) as doc:
                    page_count = len(doc)
                    text = " ".join(page.get_text() for page in doc).strip()
                if len(text) > 100:
                    return text, page_count
            except Exception:
                pass

        images = convert_from_path(file_path)
        page_count = len(images)
        ocr_text_parts = [pytesseract.image_to_string(image) for image in images]
        text = " ".join(ocr_text_parts).strip()
        if not text:
            raise RuntimeError(f"All PDF processing methods failed for file: {os.path.basename(file_path)}")
        return text, page_count

    def _process_docx(self, file_path: str) -> tuple[str, Optional[int]]:
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text), None

    def _process_txt(self, file_path: str) -> tuple[str, Optional[int]]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as handle:
            return handle.read(), None

    def _process_pptx(self, file_path: str) -> tuple[str, int]:
        prs = Presentation(file_path)
        text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
        return text, len(prs.slides)

    def _create_metadata(
        self, doc_id: str, text: str, file_path: str, page_count: Optional[int]
    ) -> DocumentMetadata:
        return DocumentMetadata(
            document_id=doc_id,
            document_type=self._detect_document_type(text),
            upload_timestamp=datetime.now().isoformat(),
            file_size=os.path.getsize(file_path),
            page_count=page_count,
            company_name=self._extract_company_name(text),
        )

    def _create_chunks(self, text: str, metadata: DocumentMetadata) -> List[Dict[str, Any]]:
        if not text:
            return []

        sentences = sent_tokenize(text)
        chunks: List[Dict[str, Any]] = []
        current_chunk_sentences: List[str] = []
        current_length = 0
        chunk_id = 0

        for sentence in sentences:
            sentence_length = len(sentence)
            if current_length + sentence_length > self.chunk_size and current_chunk_sentences:
                chunk_text = " ".join(current_chunk_sentences)
                chunks.append(self._create_chunk_dict(chunk_text, chunk_id, metadata))
                chunk_id += 1

                overlap_sentences: List[str] = []
                overlap_len = 0
                for existing in reversed(current_chunk_sentences):
                    if overlap_len + len(existing) < self.chunk_overlap:
                        overlap_sentences.insert(0, existing)
                        overlap_len += len(existing)
                    else:
                        break

                current_chunk_sentences = overlap_sentences + [sentence]
                current_length = overlap_len + sentence_length
            else:
                current_chunk_sentences.append(sentence)
                current_length += sentence_length

        if current_chunk_sentences:
            chunk_text = " ".join(current_chunk_sentences)
            chunks.append(self._create_chunk_dict(chunk_text, chunk_id, metadata))

        return chunks

    def _create_chunk_dict(
        self, text: str, chunk_id: int, metadata: DocumentMetadata
    ) -> Dict[str, Any]:
        chunk_data: Dict[str, Any] = {
            "id": f"{metadata.document_id}_{chunk_id}",
            "chunk_text": text,
            "document_id": metadata.document_id,
            "chunk_id": chunk_id,
            "document_type": metadata.document_type.value,
        }
        if metadata.company_name:
            chunk_data["company_name"] = metadata.company_name
        if metadata.page_count is not None:
            chunk_data["page_count"] = metadata.page_count
        return chunk_data

    def _clean_text(self, text: str) -> str:
        return re.sub(r"\s+", " ", text).strip()

    def _detect_document_type(self, text: str) -> DocumentType:
        return DocumentType.UNKNOWN

    def _extract_company_name(self, text: str) -> Optional[str]:
        return None

