# document_processor.py

import os
import re
from datetime import datetime
from typing import List, Dict, Any, Optional

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    print("Warning: PyMuPDF not available.")

from docx import Document
from pptx import Presentation
from nltk.tokenize import sent_tokenize

# This block ensures the necessary NLTK data is downloaded once
try:
    import nltk
    nltk.data.find('tokenizers/punkt')
    nltk.data.find('tokenizers/punkt_tab')
except (ImportError, LookupError):
    print("NLTK data not found. Downloading required packages...")
    nltk.download('punkt', quiet=True)
    nltk.download('punkt_tab', quiet=True)

import pytesseract
from pdf2image import convert_from_path

from config import Config
from models import DocumentMetadata, DocumentType


class DocumentProcessor:
    """Handles document parsing, text extraction, chunking, and metadata generation."""

    def __init__(self):
        """Initializes the processor with settings from the config."""
        self.config = Config()
        self.chunk_size = self.config.CHUNK_SIZE
        self.chunk_overlap = self.config.CHUNK_OVERLAP


    def process_document(self, file_path: str, document_id: str) -> Dict[str, Any]:
        """
        Main method to process an uploaded document.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document not found at path: {file_path}")

        text_content, metadata = self._extract_text_and_metadata(file_path, document_id)
        print(f"[DocumentProcessor] Extracted text length: {len(text_content)} for doc_id: {document_id}")

        chunks = self._create_chunks(text_content, metadata)
        print(f"[DocumentProcessor] Created {len(chunks)} chunks for doc_id: {document_id}")

        # Include 'total_chunks' for consistency with test scripts
        return {
            "document_id": document_id,
            "chunks": chunks,
            "metadata": metadata,
            "total_chunks": len(chunks)
        }

    def _extract_text_and_metadata(self, file_path: str, document_id: str) -> tuple[str, DocumentMetadata]:
        """Routes file to the correct text extraction method based on its extension."""
        ext = os.path.splitext(file_path)[1].lower()
        processor_map = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.txt': self._process_txt,
            '.pptx': self._process_pptx
        }
        if ext not in processor_map:
            raise ValueError(f"Unsupported file format: {ext}")
        
        text, page_count = processor_map[ext](file_path)
        cleaned_text = self._clean_text(text)
        metadata = self._create_metadata(document_id, cleaned_text, file_path, page_count)
        
        return cleaned_text, metadata

    def _process_pdf(self, file_path: str) -> tuple[str, int]:
        """
        Extracts text from a PDF, trying standard extraction first and falling back to OCR.
        """
        text = ""
        page_count = 0
        
        if PYMUPDF_AVAILABLE:
            try:
                with fitz.open(file_path) as doc:
                    page_count = len(doc)
                    text = " ".join(page.get_text() for page in doc).strip()
                if len(text) > 100:
                    print("[DocumentProcessor] Extracted text successfully with PyMuPDF.")
                    return text, page_count
            except Exception as e:
                print(f"PyMuPDF failed: {e}. Trying OCR.")

        print("[DocumentProcessor] Using OCR fallback for PDF...")
        try:
            images = convert_from_path(file_path)
            page_count = len(images)
            ocr_text_parts = [pytesseract.image_to_string(image) for image in images]
            text = " ".join(ocr_text_parts).strip()
            
            if len(text) > 0:
                print("[DocumentProcessor] Extracted text successfully with OCR.")
                return text, page_count
            else:
                raise RuntimeError("OCR processing resulted in empty text.")
        except Exception as ocr_error:
            print(f"❌ OCR processing failed: {ocr_error}")
            raise RuntimeError(f"All PDF processing methods failed for file: {os.path.basename(file_path)}")

    def _process_docx(self, file_path: str) -> tuple[str, Optional[int]]:
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text), None

    def _process_txt(self, file_path: str) -> tuple[str, Optional[int]]:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read(), None

    def _process_pptx(self, file_path: str) -> tuple[str, int]:
        prs = Presentation(file_path)
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text)
        return text, len(prs.slides)

    def _create_metadata(self, doc_id: str, text: str, file_path: str, page_count: Optional[int]) -> DocumentMetadata:
        return DocumentMetadata(
            document_id=doc_id,
            document_type=self._detect_document_type(text),
            upload_timestamp=datetime.now().isoformat(),
            file_size=os.path.getsize(file_path),
            page_count=page_count,
            company_name=self._extract_company_name(text),
        )

# In document_processor.py

    def _create_chunks(self, text: str, metadata: DocumentMetadata) -> List[Dict[str, Any]]:
        if not text:
            return []
        
        sentences = sent_tokenize(text)
        chunks = []
        current_chunk_sentences = []
        current_length = 0
        chunk_id = 0
        
        for sentence in sentences:
            sentence_length = len(sentence)
            # Check if adding the next sentence exceeds the chunk size
            if current_length + sentence_length > self.chunk_size and current_chunk_sentences:
                # 1. Finalize and add the current chunk
                chunk_text = " ".join(current_chunk_sentences)
                chunks.append(self._create_chunk_dict(chunk_text, chunk_id, metadata))
                chunk_id += 1
                
                # 2. Create the overlap for the next chunk
                overlap_sentences = []
                overlap_len = 0
                # Work backwards from the end of the last chunk to build the overlap
                for s in reversed(current_chunk_sentences):
                    if overlap_len + len(s) < self.chunk_overlap:
                        overlap_sentences.insert(0, s)
                        overlap_len += len(s)
                    else:
                        break
                
                # 3. Start the new chunk with the overlap and the current sentence
                current_chunk_sentences = overlap_sentences + [sentence]
                current_length = overlap_len + sentence_length
            else:
                # If chunk size is not exceeded, just add the sentence
                current_chunk_sentences.append(sentence)
                current_length += sentence_length
        
        # Add the very last chunk if any sentences are still remaining
        if current_chunk_sentences:
            chunk_text = " ".join(current_chunk_sentences)
            chunks.append(self._create_chunk_dict(chunk_text, chunk_id, metadata))
            
        return chunks

# In document_processor.py

# ... (all other code remains the same) ...

    def _create_chunk_dict(self, text: str, chunk_id: int, metadata: DocumentMetadata) -> Dict[str, Any]:
        """
        Helper to create a structured dictionary for a single chunk,
        formatted for Pinecone's integrated embedding.
        """
        # Start with the essential fields
        chunk_data = {
            "id": f"{metadata.document_id}_{chunk_id}",
            "chunk_text": text, # This key MUST match the field_map in the index
            "document_id": metadata.document_id,
            "chunk_id": chunk_id,
            "document_type": metadata.document_type.value,
        }
        
        # FIX: Only add optional fields if they have a value (are not None)
        if metadata.company_name:
            chunk_data["company_name"] = metadata.company_name
        if metadata.page_count is not None:
            chunk_data["page_count"] = metadata.page_count
            
        return chunk_data
    
    

    def _clean_text(self, text: str) -> str:
        return re.sub(r'\s+', ' ', text).strip()

    def _detect_document_type(self, text: str) -> DocumentType:
        # A more sophisticated implementation could use keyword matching or a small classifier model.
        return DocumentType.UNKNOWN
        
    def _extract_company_name(self, text: str) -> Optional[str]:
        # A more sophisticated implementation could use Named Entity Recognition (NER).
        return None
    


    # In query_engine.py

